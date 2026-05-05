import fitz  # PyMuPDF
import docx
from pptx import Presentation
from pathlib import Path
from typing import Optional, List
import logging
import subprocess
import tempfile
import os

logger = logging.getLogger(__name__)

class DocumentProcessor:
    @staticmethod
    def extract_pdf(file_path: str) -> dict:
        """Extract text and images from PDF."""
        doc = fitz.open(file_path)
        content = []
        images = []
        num_pages = len(doc)

        for page_num, page in enumerate(doc):
            text = page.get_text()
            content.append({
                "page": page_num + 1,
                "text": text,
            })

            # Extract images
            for img_idx, img in enumerate(page.get_images()):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    image_path = f"{file_path}_page{page_num+1}_img{img_idx}.png"
                    pix.save(image_path)
                    images.append({
                        "page": page_num + 1,
                        "path": image_path,
                    })
                except Exception as e:
                    logger.debug(f"Could not extract image: {e}")

        full_text = "\n".join([c["text"] for c in content])
        doc.close()

        return {
            "type": "pdf",
            "pages": num_pages,
            "content": content,
            "full_text": full_text,
            "images": images,
        }

    @staticmethod
    def extract_docx(file_path: str) -> dict:
        """Extract text from DOCX."""
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        full_text = "\n".join(paragraphs)

        return {
            "type": "docx",
            "paragraphs": len(paragraphs),
            "content": paragraphs,
            "full_text": full_text,
        }

    @staticmethod
    def convert_pptx_to_pdf(file_path: str) -> str:
        """Convert PPTX to PDF by rendering slides as images."""
        try:
            output_path = file_path.replace('.pptx', '.pdf')

            # Try LibreOffice command line first
            try:
                import subprocess
                import shutil

                # Check if soffice exists
                if shutil.which('soffice') or shutil.which('libreoffice'):
                    cmd = shutil.which('soffice') or shutil.which('libreoffice')
                    subprocess.run([
                        cmd,
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(Path(file_path).parent),
                        file_path
                    ], check=True, capture_output=True, timeout=120)

                    if os.path.exists(output_path):
                        logger.info(f"Successfully converted to PDF using LibreOffice: {output_path}")
                        return output_path
            except Exception as e:
                logger.info(f"LibreOffice not available: {e}")

            # Fallback: Convert each slide to image and embed in PDF
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from PIL import Image, ImageDraw
            import tempfile

            prs = Presentation(file_path)
            slide_width_inches = prs.slide_width / 914400
            slide_height_inches = prs.slide_height / 914400

            # Create PDF with images
            c = canvas.Canvas(output_path, pagesize=(
                slide_width_inches * 72,
                slide_height_inches * 72
            ))

            for slide_num, slide in enumerate(prs.slides):
                # Create image for this slide
                img_width, img_height = 1280, 720
                img = Image.new('RGB', (img_width, img_height), color='white')
                draw = ImageDraw.Draw(img)

                # Draw slide content
                y_offset = 40
                draw.text((40, y_offset), f"Slide {slide_num + 1}", fill='black')
                y_offset += 40

                # Extract and draw text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        try:
                            text = shape.text.encode('utf-8', errors='replace').decode('utf-8')
                            for line in text.split('\n'):
                                if y_offset < img_height - 40:
                                    draw.text((40, y_offset), line[:100], fill='black')
                                    y_offset += 20
                        except:
                            pass

                # Save image to temp file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    img.save(tmp.name)

                    # Add image to PDF
                    c.drawImage(tmp.name, 0, 0,
                               width=slide_width_inches * 72,
                               height=slide_height_inches * 72)
                    c.showPage()

                    # Clean up temp file
                    try:
                        os.unlink(tmp.name)
                    except:
                        pass

            c.save()
            logger.info(f"Successfully created PDF from images: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Error converting PPTX to PDF: {e}", exc_info=True)
            raise

    @staticmethod
    def extract_pptx(file_path: str) -> dict:
        """Extract basic info from PPTX."""
        try:
            prs = Presentation(file_path)
            full_text = ""

            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        try:
                            text = shape.text.encode('utf-8', errors='replace').decode('utf-8')
                            full_text += text + "\n"
                        except:
                            pass

            return {
                "type": "pptx",
                "slides": len(prs.slides),
                "content": [],
                "full_text": full_text,
            }
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            return {
                "type": "pptx",
                "slides": 0,
                "content": [],
                "full_text": "",
            }

    @staticmethod
    def extract_code(file_path: str) -> dict:
        """Extract code file content."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        return {
            "type": "code",
            "language": Path(file_path).suffix[1:],
            "lines": len(content.split('\n')),
            "full_text": content,
        }

    @staticmethod
    def extract_text(file_path: str) -> dict:
        """Extract plain text file content."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        return {
            "type": "text",
            "full_text": content,
        }

    @classmethod
    def process_document(cls, file_path: str) -> dict:
        """Auto-detect file type and extract content."""
        path = Path(file_path)
        ext = path.suffix.lower()

        try:
            if ext == ".pdf":
                return cls.extract_pdf(file_path)
            elif ext == ".docx":
                return cls.extract_docx(file_path)
            elif ext == ".pptx":
                return cls.extract_pptx(file_path)
            elif ext in [".py", ".js", ".ts", ".java", ".cpp", ".c", ".go", ".rs"]:
                return cls.extract_code(file_path)
            elif ext in [".txt", ".md", ".rst"]:
                return cls.extract_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            raise
